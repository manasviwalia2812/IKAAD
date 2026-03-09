from pathlib import Path
from typing import List
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document


def _load_docx(file_path: Path) -> List[Document]:
    """Load a .docx file and return LangChain Documents."""
    try:
        from docx import Document as DocxDocument
    except ImportError:
        raise ImportError("python-docx is required for DOCX support. pip install python-docx")
    doc = DocxDocument(str(file_path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    parts.append(cell.text)
    if not parts:
        return [Document(page_content="[No text content extracted]", metadata={"source": str(file_path)})]
    text = "\n\n".join(parts)
    return [Document(page_content=text, metadata={"source": str(file_path)})]


def _load_pptx(file_path: Path) -> List[Document]:
    """Load a .pptx file and return LangChain Documents (one per slide or one combined)."""
    try:
        from pptx import Presentation
        from pptx.util import Inches
    except ImportError:
        raise ImportError("python-pptx is required for PPTX support. pip install python-pptx")
    prs = Presentation(str(file_path))
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_texts.append(shape.text)
        if slide_texts:
            parts.append(f"[Slide {i + 1}]\n" + "\n".join(slide_texts))
    if not parts:
        return [Document(page_content="[No text content extracted]", metadata={"source": str(file_path)})]
    text = "\n\n".join(parts)
    return [Document(page_content=text, metadata={"source": str(file_path)})]


def load_documents_from_directory(base_dir: str) -> List[Document]:
    """
    Load all PDF, DOCX, and PPTX documents from a given directory (recursively).

    Args:
        base_dir (str): Path to the base directory containing PDFs, DOCX, and/or PPTX files.

    Returns:
        List[Document]: List of loaded LangChain Document objects.
    """
    base_path = Path(base_dir).resolve()
    documents: List[Document] = []

    if not base_path.exists():
        raise FileNotFoundError(f"Directory does not exist: {base_path}")

    # PDFs
    for pdf_file in base_path.glob("**/*.pdf"):
        loader = PyPDFLoader(str(pdf_file))
        docs = loader.load()
        for doc in docs:
            doc.metadata["source"] = str(pdf_file)
        documents.extend(docs)

    # DOCX
    for docx_file in base_path.glob("**/*.docx"):
        if docx_file.name.startswith("~$"):
            continue
        docs = _load_docx(docx_file)
        documents.extend(docs)

    # PPTX
    for pptx_file in base_path.glob("**/*.pptx"):
        if pptx_file.name.startswith("~$"):
            continue
        docs = _load_pptx(pptx_file)
        documents.extend(docs)

    return documents
